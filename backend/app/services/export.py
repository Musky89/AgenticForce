"""Document Export — generate client-ready PDFs and PPTX from agent outputs.

No client ever sees the platform. They see polished strategy decks,
concept presentations, and brand documents.
"""
import io
import logging
from pathlib import Path
from datetime import datetime

from reportlab.lib.pagesizes import A4
from reportlab.lib.units import inch, cm
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.colors import HexColor
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, PageBreak
from reportlab.lib.enums import TA_LEFT, TA_CENTER

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

logger = logging.getLogger(__name__)

EXPORTS_DIR = Path(__file__).parent.parent.parent / "exports"
EXPORTS_DIR.mkdir(exist_ok=True)

BRAND_DARK = HexColor("#0a0a0f")
BRAND_ACCENT = HexColor("#6366f1")
BRAND_TEXT = HexColor("#e4e4e7")


def _build_pdf_styles():
    styles = getSampleStyleSheet()
    styles.add(ParagraphStyle("DocTitle", fontSize=28, textColor=BRAND_ACCENT, spaceAfter=20, fontName="Helvetica-Bold"))
    styles.add(ParagraphStyle("DocSubtitle", fontSize=14, textColor=BRAND_TEXT, spaceAfter=30, fontName="Helvetica"))
    styles.add(ParagraphStyle("SectionHead", fontSize=18, textColor=BRAND_ACCENT, spaceBefore=20, spaceAfter=10, fontName="Helvetica-Bold"))
    styles.add(ParagraphStyle("BodyText2", fontSize=11, textColor=HexColor("#d4d4d8"), leading=16, spaceAfter=8, fontName="Helvetica"))
    styles.add(ParagraphStyle("MetaText", fontSize=9, textColor=HexColor("#71717a"), spaceAfter=4, fontName="Helvetica"))
    return styles


def export_strategy_pdf(
    client_name: str,
    project_name: str,
    strategy_content: str,
    research_content: str | None = None,
) -> str:
    """Export strategy document as a polished PDF."""
    filename = f"strategy_{project_name.replace(' ', '_').lower()}_{datetime.utcnow().strftime('%Y%m%d')}.pdf"
    filepath = EXPORTS_DIR / filename

    doc = SimpleDocTemplate(str(filepath), pagesize=A4, topMargin=2*cm, bottomMargin=2*cm, leftMargin=2.5*cm, rightMargin=2.5*cm)
    styles = _build_pdf_styles()
    story = []

    story.append(Paragraph(f"{client_name}", styles["DocTitle"]))
    story.append(Paragraph(f"{project_name} — Creative Strategy", styles["DocSubtitle"]))
    story.append(Paragraph(f"Prepared by AgenticForce | {datetime.utcnow().strftime('%B %Y')}", styles["MetaText"]))
    story.append(Spacer(1, 30))

    if research_content:
        story.append(Paragraph("Research & Insights", styles["SectionHead"]))
        for line in research_content.split("\n"):
            if line.strip():
                story.append(Paragraph(line.strip(), styles["BodyText2"]))
        story.append(PageBreak())

    story.append(Paragraph("Creative Strategy", styles["SectionHead"]))
    for line in strategy_content.split("\n"):
        line = line.strip()
        if not line:
            story.append(Spacer(1, 6))
        elif line.startswith("**") and line.endswith("**"):
            story.append(Paragraph(line.strip("*"), styles["SectionHead"]))
        elif line.startswith("#"):
            story.append(Paragraph(line.lstrip("# "), styles["SectionHead"]))
        else:
            story.append(Paragraph(line, styles["BodyText2"]))

    doc.build(story)
    logger.info(f"Exported strategy PDF: {filename}")
    return filename


def export_concept_pptx(
    client_name: str,
    project_name: str,
    concepts_content: str,
    art_direction_content: str | None = None,
) -> str:
    """Export concept presentation as PPTX."""
    filename = f"concepts_{project_name.replace(' ', '_').lower()}_{datetime.utcnow().strftime('%Y%m%d')}.pptx"
    filepath = EXPORTS_DIR / filename

    prs = Presentation()
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)

    # Title slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(0x0a, 0x0a, 0x0f)

    title_box = slide.shapes.add_textbox(Inches(1.5), Inches(3), Inches(13), Inches(2))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = client_name
    p.font.size = Pt(48)
    p.font.color.rgb = RGBColor(0x63, 0x66, 0xf1)
    p.font.bold = True
    p.alignment = PP_ALIGN.LEFT

    p2 = tf.add_paragraph()
    p2.text = f"{project_name} — Creative Concepts"
    p2.font.size = Pt(24)
    p2.font.color.rgb = RGBColor(0xe4, 0xe4, 0xe7)

    p3 = tf.add_paragraph()
    p3.text = f"Prepared by AgenticForce | {datetime.utcnow().strftime('%B %Y')}"
    p3.font.size = Pt(14)
    p3.font.color.rgb = RGBColor(0x71, 0x71, 0x7a)

    # Content slides
    sections = concepts_content.split("\n\n")
    current_section = []

    for section in sections:
        lines = section.strip().split("\n")
        if not lines or not lines[0].strip():
            continue

        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = RGBColor(0x0a, 0x0a, 0x0f)

        content_box = slide.shapes.add_textbox(Inches(1.5), Inches(1), Inches(13), Inches(7))
        tf = content_box.text_frame
        tf.word_wrap = True

        for i, line in enumerate(lines):
            line = line.strip()
            if not line:
                continue
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()

            if line.startswith("**") or line.startswith("#"):
                p.text = line.strip("*# ")
                p.font.size = Pt(28)
                p.font.color.rgb = RGBColor(0x63, 0x66, 0xf1)
                p.font.bold = True
            else:
                p.text = line
                p.font.size = Pt(16)
                p.font.color.rgb = RGBColor(0xd4, 0xd4, 0xd8)
            p.space_after = Pt(8)

    # Art direction slide(s)
    if art_direction_content:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = RGBColor(0x0a, 0x0a, 0x0f)

        header = slide.shapes.add_textbox(Inches(1.5), Inches(1), Inches(13), Inches(1))
        hp = header.text_frame.paragraphs[0]
        hp.text = "Art Direction"
        hp.font.size = Pt(36)
        hp.font.color.rgb = RGBColor(0x63, 0x66, 0xf1)
        hp.font.bold = True

        body = slide.shapes.add_textbox(Inches(1.5), Inches(2.5), Inches(13), Inches(5.5))
        btf = body.text_frame
        btf.word_wrap = True
        for line in art_direction_content.split("\n")[:30]:
            p = btf.add_paragraph()
            p.text = line.strip()
            p.font.size = Pt(14)
            p.font.color.rgb = RGBColor(0xd4, 0xd4, 0xd8)

    prs.save(str(filepath))
    logger.info(f"Exported concept PPTX: {filename}")
    return filename


def export_deliverables_pdf(
    client_name: str,
    project_name: str,
    deliverables: list[dict],
) -> str:
    """Export all deliverables as a consolidated PDF."""
    filename = f"deliverables_{project_name.replace(' ', '_').lower()}_{datetime.utcnow().strftime('%Y%m%d')}.pdf"
    filepath = EXPORTS_DIR / filename

    doc = SimpleDocTemplate(str(filepath), pagesize=A4, topMargin=2*cm, bottomMargin=2*cm, leftMargin=2.5*cm, rightMargin=2.5*cm)
    styles = _build_pdf_styles()
    story = []

    story.append(Paragraph(client_name, styles["DocTitle"]))
    story.append(Paragraph(f"{project_name} — Deliverables Package", styles["DocSubtitle"]))
    story.append(Paragraph(f"Prepared by AgenticForce | {datetime.utcnow().strftime('%B %Y')}", styles["MetaText"]))
    story.append(Spacer(1, 20))

    for d in deliverables:
        story.append(Paragraph(d.get("title", "Untitled"), styles["SectionHead"]))
        if d.get("pipeline_stage"):
            story.append(Paragraph(f"Stage: {d['pipeline_stage'].replace('_', ' ').title()}", styles["MetaText"]))
        for line in d.get("content", "").split("\n"):
            if line.strip():
                story.append(Paragraph(line.strip(), styles["BodyText2"]))
        story.append(PageBreak())

    doc.build(story)
    logger.info(f"Exported deliverables PDF: {filename}")
    return filename
